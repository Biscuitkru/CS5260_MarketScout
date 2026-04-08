"""
MarketScout: Slide Deck Generator
====================================
Parses the publisher's 4-slide markdown output and produces a .pptx file.
"""
import io
import json
import os
import re

from langchain_core.messages import HumanMessage, SystemMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Theme — minimalist white + pale yellow
_BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white
_TITLE   = RGBColor(0x1A, 0x1A, 0x1A)   # near-black
_BODY    = RGBColor(0x33, 0x33, 0x33)   # dark gray
_SUB     = RGBColor(0x88, 0x88, 0x88)   # medium gray
_RULE    = RGBColor(0xF5, 0xE0, 0x7A)   # pale yellow rule line


def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _BG


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _parse_inline(text: str) -> list[tuple[str, bool]]:
    """
    Parse **bold** markdown into a list of (text, is_bold) tuples.
    Strips all ** markers so PowerPoint never sees raw asterisks.
    """
    parts = []
    for segment in re.split(r"(\*\*[^*]+\*\*)", text):
        if segment.startswith("**") and segment.endswith("**"):
            parts.append((segment[2:-2], True))
        elif segment:
            parts.append((segment, False))
    return parts


def _add_rich_para(tf, text: str, font_size: int, color: RGBColor,
                   space_before: int = 0, first: bool = False):
    """Add a paragraph with inline bold support, no raw markdown asterisks."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if space_before:
        p.space_before = Pt(space_before)

    for segment, bold in _parse_inline(text):
        run = p.add_run()
        run.text = segment
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def summarize_to_slides(report: str) -> list[dict]:
    """
    Ask the LLM to condense a full markdown report into 4 slides.
    Returns a list of dicts: [{title, subtitle, bullets}, ...]
    """
    system = """\
You are a presentation designer. Given a market research report, condense it into exactly 4 slides.

Output valid JSON only — no markdown fences, no extra text.

Format:
[
  {
    "title": "Market Overview",
    "subtitle": "one-line context",
    "bullets": ["point 1", "point 2", "point 3"]
  },
  ...
]

Rules:
- Exactly 4 slides in this order:
    1. Market Overview
    2. Competitive Landscape
    3. Customer Insights & Gaps
    4. Recommendations & Risks
- 3 to 4 bullets per slide maximum
- Each bullet is one tight sentence — no markdown, no bold asterisks, no filler
- Bullets must be specific and actionable, not generic summaries
"""
    model_name = os.getenv("PUBLISHER_MODEL", "gemini-2.5-flash")
    llm = ChatGoogleGenerativeAI(model=model_name, temperature=0.2)
    response = llm.invoke([
        SystemMessage(content=system),
        HumanMessage(content=report),
    ])

    raw = response.content.strip()
    # Strip accidental markdown code fences if present
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    slides = json.loads(raw)
    return slides[:4]


def build_pptx(slides_data: list[dict]) -> bytes:
    """Convert pre-computed slide data into a .pptx bytes object."""

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    total = len(slides_data)

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide)

        # Top accent bar (full width, thin)
        _add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.07), _RULE)

        # Bottom accent bar
        _add_rect(slide, Inches(0), Inches(7.43), Inches(13.33), Inches(0.07), _RULE)

        # Slide number — bottom right
        num_box = slide.shapes.add_textbox(Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.3))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_run = num_p.add_run()
        num_run.text = f"{i + 1}  /  {total}"
        num_run.font.size = Pt(11)
        num_run.font.color.rgb = _SUB

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.8), Inches(1.1))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_run = title_p.add_run()
        title_run.text = sd["title"]
        title_run.font.size = Pt(38)
        title_run.font.bold = True
        title_run.font.color.rgb = _TITLE

        # Subtitle / context line
        if sd["subtitle"]:
            sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(11.8), Inches(0.5))
            sub_tf = sub_box.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_run = sub_p.add_run()
            sub_run.text = sd["subtitle"]
            sub_run.font.size = Pt(15)
            sub_run.font.italic = True
            sub_run.font.color.rgb = _SUB

        # Divider rule below title area
        rule_top = Inches(1.85) if sd["subtitle"] else Inches(1.5)
        _add_rect(slide, Inches(0.6), rule_top, Inches(11.8), Inches(0.025), _RULE)

        # Bullets
        if sd["bullets"]:
            bullet_top = rule_top + Inches(0.25)
            bx = slide.shapes.add_textbox(Inches(0.6), bullet_top, Inches(11.8), Inches(5.0))
            btf = bx.text_frame
            btf.word_wrap = True

            for j, bullet in enumerate(sd["bullets"]):
                _add_rich_para(
                    btf,
                    f"\u2013   {bullet}",   # en-dash bullet
                    font_size=22,
                    color=_BODY,
                    space_before=14,
                    first=(j == 0),
                )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
